"""Generate final project presentation for OddsWatch."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Color palette ─────────────────────────────────────────────
NAVY = RGBColor(25, 60, 120)
DARK_TEXT = RGBColor(40, 40, 40)
GRAY_TEXT = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(240, 244, 250)
ACCENT_ORANGE = RGBColor(230, 140, 50)
ACCENT_GREEN = RGBColor(60, 160, 100)
ACCENT_BLUE = RGBColor(70, 130, 200)
BRONZE_COLOR = RGBColor(205, 170, 125)
SILVER_COLOR = RGBColor(180, 190, 210)
GOLD_COLOR = RGBColor(220, 190, 80)


def add_bg(slide, color=LIGHT_BG):
    """Add a solid background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text, y=Inches(0.3)):
    """Add a navy title bar across the top."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y, Inches(13.33), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_text_box(slide, left, top, width, height, text, font_size=16,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=14):
    """Add a bulleted list to a slide."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_box_shape(slide, left, top, width, height, text, fill_color,
                  font_size=11, text_color=DARK_TEXT):
    """Add a colored rectangle with centered text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = NAVY
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, left, top, width, height):
    """Add a downward arrow."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY_TEXT
    arrow.line.fill.background()
    return arrow


def add_right_arrow(slide, left, top, width, height):
    """Add a right arrow."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY_TEXT
    arrow.line.fill.background()
    return arrow


def build_presentation():
    """Build the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1: Title ────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, NAVY)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "OddsWatch", font_size=54, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "Sports Betting Odds Analytics Platform",
                 font_size=22, color=RGBColor(180, 200, 230),
                 alignment=PP_ALIGN.CENTER)

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(4.5), Inches(4.0), Inches(4), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()

    add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
                 "MGMT 59000  |  Cloud Data Engineering  |  Final Project",
                 font_size=16, color=RGBColor(160, 180, 210),
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
                 "Team 2: Danny Affleck, Aadi Gupta, Ken Nanayama  |  July 2026",
                 font_size=16, color=RGBColor(160, 180, 210),
                 alignment=PP_ALIGN.CENTER)

    # ── SLIDE 2: Agenda ───────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "Agenda")

    items = [
        "1.  Project Overview & Problem Statement",
        "2.  Architecture & Pattern Justification",
        "3.  Data Model & Star Schema",
        "4.  SCD Strategy",
        "5.  AWS vs. GCP Tool Selection",
        "6.  Partial Deployment: dbt on Athena",
        "7.  Cloud Cost Estimate",
        "8.  Scale Analysis: What Breaks at 10x",
        "9.  Trade-Off Analysis: Dominant Dimension",
        "10. Live Demo & Q&A",
    ]
    add_bullet_list(slide, Inches(1), Inches(1.6), Inches(10), Inches(5.5),
                    items, font_size=18)

    # ── SLIDE 3: Problem Statement ────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "Project Overview")

    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.8),
                 "What is OddsWatch?", font_size=22, bold=True, color=NAVY)
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.2),
                 "A cloud-native analytics platform that ingests sports data from "
                 "multiple public sources, normalizes it into a unified schema, "
                 "generates synthetic betting odds, and serves it as a star schema "
                 "optimized for analytical queries via AWS Athena.",
                 font_size=16)

    add_text_box(slide, Inches(0.8), Inches(3.8), Inches(5), Inches(0.5),
                 "Data Sources", font_size=18, bold=True, color=NAVY)
    add_bullet_list(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(1.5),
                    ["FiveThirtyEight MLB ELO (~230K games)",
                     "FIFA World Cup Results (~1K matches)"],
                    font_size=14)

    add_text_box(slide, Inches(7), Inches(3.8), Inches(5), Inches(0.5),
                 "Key Capabilities", font_size=18, bold=True, color=NAVY)
    add_bullet_list(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(2.5),
                    ["Batch pipeline: ingest, transform, serve",
                     "Medallion architecture (Bronze/Silver/Gold)",
                     "Star schema with 4 dims + 1 fact table",
                     "Synthetic odds: spread, total, moneyline",
                     "Glue Catalog + Athena SQL access"],
                    font_size=14)

    # ── SLIDE 4: Architecture Diagram ─────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "Architecture: Medallion Lakehouse")

    # Layout: 3 vertical swim lanes (Sources | Pipeline | Storage & Serving)
    # with clear left-to-right flow per row

    # Lane labels
    add_text_box(slide, Inches(0.4), Inches(1.35), Inches(3), Inches(0.4),
                 "DATA SOURCES", font_size=11, bold=True, color=GRAY_TEXT)
    add_text_box(slide, Inches(4.0), Inches(1.35), Inches(3), Inches(0.4),
                 "PIPELINE", font_size=11, bold=True, color=GRAY_TEXT)
    add_text_box(slide, Inches(8.2), Inches(1.35), Inches(3), Inches(0.4),
                 "STORAGE", font_size=11, bold=True, color=GRAY_TEXT)
    add_text_box(slide, Inches(11.2), Inches(1.35), Inches(2), Inches(0.4),
                 "SERVING", font_size=11, bold=True, color=GRAY_TEXT)

    # Row 1: Bronze layer
    brow_y = Inches(1.9)
    add_box_shape(slide, Inches(0.4), brow_y, Inches(1.6), Inches(0.65),
                  "FiveThirtyEight\nMLB ELO", ACCENT_ORANGE, font_size=9, text_color=WHITE)
    add_box_shape(slide, Inches(2.15), brow_y, Inches(1.6), Inches(0.65),
                  "GitHub\nFIFA World Cup", ACCENT_ORANGE, font_size=9, text_color=WHITE)
    add_right_arrow(slide, Inches(3.85), brow_y + Inches(0.12), Inches(0.4), Inches(0.35))
    add_box_shape(slide, Inches(4.35), brow_y, Inches(3.5), Inches(0.65),
                  "Ingest (httpx)\nDownload + Upload to S3", RGBColor(255, 220, 200), font_size=10)
    add_right_arrow(slide, Inches(7.95), brow_y + Inches(0.12), Inches(0.4), Inches(0.35))
    add_box_shape(slide, Inches(8.45), brow_y, Inches(2.5), Inches(0.65),
                  "S3: bronze/\nCSV (raw)", BRONZE_COLOR, font_size=10)

    # Bronze label
    bronze_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(11.3), brow_y, Inches(1.5), Inches(0.65))
    bronze_label.fill.solid()
    bronze_label.fill.fore_color.rgb = BRONZE_COLOR
    bronze_label.line.fill.background()
    tf = bronze_label.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "BRONZE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Arrow down between rows
    add_arrow(slide, Inches(6.0), Inches(2.6), Inches(0.3), Inches(0.4))

    # Row 2: Silver layer
    srow_y = Inches(3.15)
    add_box_shape(slide, Inches(4.35), srow_y, Inches(3.5), Inches(0.9),
                  "bronze_to_silver Transform\nNormalize schemas + Generate\nsynthetic odds (seeded)",
                  RGBColor(200, 210, 240), font_size=10)
    add_right_arrow(slide, Inches(7.95), srow_y + Inches(0.25), Inches(0.4), Inches(0.35))
    add_box_shape(slide, Inches(8.45), srow_y, Inches(2.5), Inches(0.9),
                  "S3: silver/\nParquet (unified\nSilverGame schema)",
                  SILVER_COLOR, font_size=10)

    silver_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(11.3), srow_y + Inches(0.12), Inches(1.5), Inches(0.65))
    silver_label.fill.solid()
    silver_label.fill.fore_color.rgb = SILVER_COLOR
    silver_label.line.fill.background()
    tf = silver_label.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "SILVER"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Arrow down
    add_arrow(slide, Inches(6.0), Inches(4.1), Inches(0.3), Inches(0.4))

    # Row 3: Gold layer
    grow_y = Inches(4.65)
    add_box_shape(slide, Inches(4.35), grow_y, Inches(3.5), Inches(0.9),
                  "silver_to_gold Transform\nBuild star schema:\n4 dims + 1 fact table",
                  RGBColor(200, 210, 240), font_size=10)
    add_right_arrow(slide, Inches(7.95), grow_y + Inches(0.25), Inches(0.4), Inches(0.35))
    add_box_shape(slide, Inches(8.45), grow_y, Inches(2.5), Inches(0.9),
                  "S3: gold/\nParquet (star schema\n5 tables)", GOLD_COLOR, font_size=10)

    gold_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(11.3), grow_y + Inches(0.12), Inches(1.5), Inches(0.65))
    gold_label.fill.solid()
    gold_label.fill.fore_color.rgb = GOLD_COLOR
    gold_label.line.fill.background()
    tf = gold_label.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "GOLD"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # Arrow down to serving layer
    add_arrow(slide, Inches(9.55), Inches(5.6), Inches(0.3), Inches(0.4))

    # Row 4: Serving layer
    serve_y = Inches(6.15)
    add_box_shape(slide, Inches(8.45), serve_y, Inches(2.5), Inches(0.65),
                  "AWS Glue\nData Catalog", RGBColor(210, 190, 230), font_size=10)
    add_right_arrow(slide, Inches(11.05), serve_y + Inches(0.12), Inches(0.4), Inches(0.35))
    add_box_shape(slide, Inches(11.55), serve_y, Inches(1.4), Inches(0.65),
                  "Athena\nSQL Queries", RGBColor(210, 190, 230), font_size=10)

    # dbt integration (separate column on far left, connected to gold)
    dbt_y = Inches(4.65)
    add_box_shape(slide, Inches(0.4), dbt_y, Inches(1.6), Inches(0.9),
                  "dbt Project\n7 models\n45 tests", ACCENT_GREEN, font_size=10, text_color=WHITE)
    add_right_arrow(slide, Inches(2.1), dbt_y + Inches(0.25), Inches(0.5), Inches(0.35))
    add_text_box(slide, Inches(2.6), dbt_y + Inches(0.15), Inches(1.5), Inches(0.5),
                 "reads from\nGlue Catalog", font_size=9, color=GRAY_TEXT)

    # Orchestrator box
    add_box_shape(slide, Inches(0.4), Inches(6.15), Inches(3.3), Inches(0.65),
                  "Pipeline Orchestrator (pipeline.py)", RGBColor(240, 240, 220), font_size=10)

    # Pattern label at bottom
    add_text_box(slide, Inches(4.35), Inches(6.95), Inches(8), Inches(0.4),
                 "Pattern: Medallion Lakehouse  |  Format: CSV (bronze) / Parquet (silver, gold)  |  Catalog: AWS Glue",
                 font_size=11, bold=True, color=NAVY)

    # ── SLIDE 5: Why Lakehouse ────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "Why Lakehouse Over Alternatives?")

    # Three columns
    for i, (title, desc, verdict) in enumerate([
        ("Lambda Architecture",
         "Separate batch + speed layers.\nMaintain two codepaths.",
         "Rejected: batch-only today.\nNo benefit until streaming added."),
        ("Medallion Lakehouse",
         "Single S3 lake with progressive\nrefinement (Bronze/Silver/Gold).\nGlue Catalog for schema.",
         "Chosen: simple, extensible,\nlow cost, single source of truth."),
        ("Data Mesh",
         "Domain-owned data products.\nFederated governance.",
         "Rejected: premature for\na 2-source, single-team project."),
    ]):
        left = Inches(0.5 + i * 4.2)
        color = ACCENT_GREEN if i == 1 else RGBColor(200, 200, 200)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left, Inches(1.8), Inches(3.8), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(3)

        add_text_box(slide, left + Inches(0.2), Inches(2.0), Inches(3.4), Inches(0.5),
                     title, font_size=18, bold=True, color=NAVY)
        add_text_box(slide, left + Inches(0.2), Inches(2.7), Inches(3.4), Inches(1.5),
                     desc, font_size=13, color=DARK_TEXT)
        verdict_color = ACCENT_GREEN if i == 1 else RGBColor(180, 60, 60)
        add_text_box(slide, left + Inches(0.2), Inches(4.5), Inches(3.4), Inches(1.5),
                     verdict, font_size=13, bold=True, color=verdict_color)

    # ── SLIDE 6: Data Model - Star Schema ─────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "Gold Layer: Star Schema")

    # Central fact table - large, prominent
    fact_x, fact_y = Inches(4.2), Inches(2.4)
    fact_w, fact_h = Inches(5), Inches(3.8)
    fact = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fact_x, fact_y, fact_w, fact_h)
    fact.fill.solid()
    fact.fill.fore_color.rgb = RGBColor(255, 243, 215)
    fact.line.color.rgb = ACCENT_ORANGE
    fact.line.width = Pt(2.5)
    tf = fact.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = Inches(0.15)
    tf.margin_left = Inches(0.2)
    # Title
    p = tf.paragraphs[0]
    p.text = "FactGameOdds"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE
    p.alignment = PP_ALIGN.CENTER
    # Separator
    p = tf.add_paragraph()
    p.text = ""
    p.space_before = Pt(4)
    # FK columns
    for col_text in ["game_key  (FK)", "date_key  (FK)",
                     "home_team_key  (FK)", "away_team_key  (FK)"]:
        p = tf.add_paragraph()
        p.text = col_text
        p.font.size = Pt(11)
        p.font.color.rgb = ACCENT_BLUE
        p.font.bold = True
        p.space_before = Pt(1)
    # Measure columns
    p = tf.add_paragraph()
    p.text = ""
    p.space_before = Pt(4)
    for col_text in ["sport  |  closing_spread  |  closing_total",
                     "moneyline_home  |  moneyline_away",
                     "home_score  |  away_score",
                     "cover (bool)  |  over_under_result"]:
        p = tf.add_paragraph()
        p.text = col_text
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(1)

    # Helper to draw a dimension box
    dim_blue = RGBColor(220, 232, 250)
    dim_border = ACCENT_BLUE

    def add_dim_box(x, y, w, h, title, columns):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = dim_blue
        shape.line.color.rgb = dim_border
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.1)
        tf.margin_left = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER
        for col in columns:
            p = tf.add_paragraph()
            p.text = col
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(1)

    # DimTeam - top left
    add_dim_box(Inches(0.3), Inches(1.5), Inches(3.2), Inches(1.6),
                "DimTeam",
                ["team_key (PK)", "team_id  |  team_name", "sport"])
    # Connector: right arrow from DimTeam to Fact
    add_right_arrow(slide, Inches(3.55), Inches(2.15), Inches(0.55), Inches(0.3))

    # DimDate - top right
    add_dim_box(Inches(9.85), Inches(1.5), Inches(3.2), Inches(1.8),
                "DimDate",
                ["date_key (PK)", "date  |  day_of_week", "month  |  year", "is_weekend"])
    # Connector: left arrow (use right arrow shape, it points from dim to fact)
    left_arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      Inches(9.25), Inches(2.25), Inches(0.55), Inches(0.3))
    left_arr.fill.solid()
    left_arr.fill.fore_color.rgb = GRAY_TEXT
    left_arr.line.fill.background()
    left_arr.rotation = 180.0

    # DimGame - bottom left
    add_dim_box(Inches(0.3), Inches(4.5), Inches(3.2), Inches(1.8),
                "DimGame",
                ["game_key (PK)", "game_id  |  sport  |  season", "date  |  stage  |  venue",
                 "home_team_key  |  away_team_key"])
    add_right_arrow(slide, Inches(3.55), Inches(5.25), Inches(0.55), Inches(0.3))

    # DimMarket - bottom right
    add_dim_box(Inches(9.85), Inches(4.8), Inches(3.2), Inches(1.3),
                "DimMarket (static)",
                ["market_key (PK)", "spread  |  total  |  moneyline"])
    left_arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       Inches(9.25), Inches(5.3), Inches(0.55), Inches(0.3))
    left_arr2.fill.solid()
    left_arr2.fill.fore_color.rgb = GRAY_TEXT
    left_arr2.line.fill.background()
    left_arr2.rotation = 180.0

    # Label
    add_text_box(slide, Inches(0.3), Inches(6.7), Inches(12), Inches(0.4),
                 "Star schema: 4 dimension tables surround 1 central fact table  |  "
                 "13 columns in fact  |  Derived metrics: cover (bool), over_under_result",
                 font_size=11, bold=True, color=NAVY)

    # ── SLIDE 7: SCD Strategy ─────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "SCD Strategy")

    data = [
        ("DimTeam", "Type 1 (Overwrite)", "Team names rarely change; no need to track history"),
        ("DimGame", "Type 0 (Fixed)", "Game facts are immutable once recorded"),
        ("DimDate", "Type 0 (Fixed)", "Calendar attributes are deterministic"),
        ("DimMarket", "Type 0 (Fixed)", "Static 3-row reference table"),
        ("DimBook*", "Type 2 (History)", "Sportsbook fees change; valid_from/valid_to/is_current"),
    ]

    # Header
    add_text_box(slide, Inches(1), Inches(1.6), Inches(3), Inches(0.5),
                 "Dimension", font_size=16, bold=True, color=WHITE)
    add_text_box(slide, Inches(4), Inches(1.6), Inches(3), Inches(0.5),
                 "SCD Type", font_size=16, bold=True, color=WHITE)
    add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5),
                 "Rationale", font_size=16, bold=True, color=WHITE)
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.55))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = NAVY
    header_bar.line.fill.background()
    # Move header bar behind text
    sp = header_bar._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

    for i, (dim, scd, rationale) in enumerate(data):
        y = Inches(2.3 + i * 0.7)
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.8), y, Inches(11.7), Inches(0.6))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(235, 240, 248)
        row_bg.line.fill.background()
        sp = row_bg._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

        add_text_box(slide, Inches(1), y + Inches(0.05), Inches(3), Inches(0.5),
                     dim, font_size=14, bold=True)
        add_text_box(slide, Inches(4), y + Inches(0.05), Inches(3), Inches(0.5),
                     scd, font_size=14, color=ACCENT_BLUE)
        add_text_box(slide, Inches(7), y + Inches(0.05), Inches(5.5), Inches(0.5),
                     rationale, font_size=13)

    add_text_box(slide, Inches(1), Inches(6.0), Inches(10), Inches(0.5),
                 "* DimBook defined in model for future streaming layer; not populated in batch pipeline.",
                 font_size=12, color=GRAY_TEXT)

    # ── SLIDE 8: AWS vs GCP ───────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "AWS vs. GCP Tool Selection")

    services = [
        ("Object Storage", "Amazon S3", "Cloud Storage"),
        ("Data Catalog", "AWS Glue Catalog", "Data Catalog"),
        ("ETL / Transform", "AWS Glue Jobs", "Dataflow / Dataproc"),
        ("SQL Analytics", "Amazon Athena", "BigQuery"),
        ("Orchestration", "Step Functions", "Cloud Composer"),
        ("Compute", "AWS Lambda", "Cloud Functions"),
        ("Streaming*", "Kinesis", "Pub/Sub"),
    ]

    # Column headers
    for col, (text, x) in enumerate([
        ("Component", Inches(0.8)),
        ("AWS (Chosen)", Inches(4)),
        ("GCP Equivalent", Inches(7.5)),
    ]):
        add_text_box(slide, x, Inches(1.6), Inches(3), Inches(0.5),
                     text, font_size=15, bold=True, color=WHITE)

    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.6), Inches(1.55), Inches(10), Inches(0.55))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = NAVY
    header_bar.line.fill.background()
    sp = header_bar._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

    for i, (component, aws, gcp) in enumerate(services):
        y = Inches(2.3 + i * 0.6)
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.6), y, Inches(10), Inches(0.55))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(235, 240, 248)
        row_bg.line.fill.background()
        sp = row_bg._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

        add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(3), Inches(0.45),
                     component, font_size=13, bold=True)
        add_text_box(slide, Inches(4), y + Inches(0.05), Inches(3), Inches(0.45),
                     aws, font_size=13, color=ACCENT_BLUE)
        add_text_box(slide, Inches(7.5), y + Inches(0.05), Inches(3), Inches(0.45),
                     gcp, font_size=13)

    add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.6),
                 "Why AWS:  S3 + Glue + Athena is tightly integrated, serverless, "
                 "and aligns with course curriculum. boto3 provides mature Python SDK support.",
                 font_size=14, bold=True, color=NAVY)

    # ── SLIDE 9: Partial Deployment ───────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "Partial Deployment: dbt on Athena")

    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
                 "dbt Project Structure", font_size=20, bold=True, color=NAVY)
    add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4),
                    ["2 staging models (views) read from Glue Catalog",
                     "  - stg_mlb_games, stg_world_cup_games",
                     "5 gold models (materialized as Parquet tables)",
                     "  - dim_team, dim_game, dim_date, dim_market",
                     "  - fact_game_odds",
                     "45 data quality tests",
                     "  - not_null, unique on surrogate keys",
                     "  - FK relationships between fact and dims",
                     "  - accepted_values on sport, market_type"],
                    font_size=14)

    add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5),
                 "Demo Commands", font_size=20, bold=True, color=NAVY)

    # Command boxes
    for i, (cmd, desc) in enumerate([
        ("dbt run --profiles-dir .", "Materializes 7 models in Athena"),
        ("dbt test --profiles-dir .", "Runs 45 data quality tests"),
        ("SELECT ... FROM fact_game_odds", "Query star schema in Athena console"),
    ]):
        y = Inches(2.3 + i * 1.4)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(7), y, Inches(5.5), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(40, 44, 52)
        box.line.fill.background()
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = cmd
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT_GREEN
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        tf.margin_left = Inches(0.2)

        add_text_box(slide, Inches(7.2), y + Inches(0.65), Inches(5), Inches(0.4),
                     desc, font_size=12, color=GRAY_TEXT)

    # ── SLIDE 10: Cost Estimate ───────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "Cloud Cost Estimate")

    costs = [
        ("S3 Storage", "~20 MB", "$0.01"),
        ("S3 Requests", "~1,200/mo", "$0.02"),
        ("Glue Catalog", "12 tables", "$0.00"),
        ("Athena Queries", "~100 queries", "$0.01"),
        ("Lambda", "30 invocations", "$0.01"),
        ("Data Transfer", "~50 MB out", "$0.00"),
    ]

    for col, (text, x) in enumerate([
        ("Service", Inches(1.5)), ("Usage", Inches(5)), ("Monthly Cost", Inches(8.5)),
    ]):
        add_text_box(slide, x, Inches(1.7), Inches(3), Inches(0.5),
                     text, font_size=15, bold=True, color=WHITE)

    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(1.3), Inches(1.65), Inches(10), Inches(0.55))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = NAVY
    header_bar.line.fill.background()
    sp = header_bar._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

    for i, (service, usage, cost) in enumerate(costs):
        y = Inches(2.4 + i * 0.55)
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(1.3), y, Inches(10), Inches(0.5))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(235, 240, 248)
        row_bg.line.fill.background()
        sp = row_bg._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

        add_text_box(slide, Inches(1.5), y + Inches(0.03), Inches(3), Inches(0.45),
                     service, font_size=13)
        add_text_box(slide, Inches(5), y + Inches(0.03), Inches(3), Inches(0.45),
                     usage, font_size=13)
        add_text_box(slide, Inches(8.5), y + Inches(0.03), Inches(2), Inches(0.45),
                     cost, font_size=13, bold=True, color=ACCENT_GREEN)

    # Total
    total_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(3.5), Inches(5.8), Inches(6), Inches(0.8))
    total_box.fill.solid()
    total_box.fill.fore_color.rgb = NAVY
    total_box.line.fill.background()
    tf = total_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Estimated Total:  < $1.00 / month   (Free Tier eligible)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # ── SLIDE 11: Scale Analysis ──────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "Scale Analysis: What Breaks at 10x")

    components = [
        ("Ingestion", "100x", "Timeout / OOM", "Chunked + async httpx"),
        ("Transform", "10x", "Slow row-by-row loop", "Vectorize with numpy / PySpark"),
        ("Storage", "10x", "No partition pruning", "Hive-style partitioning"),
        ("Catalog/Query", "10x", "Full scans, slow queries", "Partition pruning + Iceberg"),
        ("Orchestration", "10x", "No retry, full re-run", "AWS Step Functions"),
    ]

    headers = ["Component", "Breaks At", "Symptom", "Mitigation"]
    x_positions = [Inches(0.5), Inches(3), Inches(5.3), Inches(8.5)]
    widths = [Inches(2.5), Inches(2.3), Inches(3.2), Inches(4)]

    for j, (header, x) in enumerate(zip(headers, x_positions)):
        add_text_box(slide, x, Inches(1.7), widths[j], Inches(0.5),
                     header, font_size=15, bold=True, color=WHITE)

    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.3), Inches(1.65), Inches(12.5), Inches(0.55))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = NAVY
    header_bar.line.fill.background()
    sp = header_bar._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

    for i, (comp, breaks, symptom, fix) in enumerate(components):
        y = Inches(2.4 + i * 0.7)
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.3), y, Inches(12.5), Inches(0.6))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(235, 240, 248)
        row_bg.line.fill.background()
        sp = row_bg._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

        add_text_box(slide, Inches(0.5), y + Inches(0.05), Inches(2.5), Inches(0.5),
                     comp, font_size=14, bold=True)
        add_text_box(slide, Inches(3), y + Inches(0.05), Inches(2.3), Inches(0.5),
                     breaks, font_size=14, color=RGBColor(180, 60, 60))
        add_text_box(slide, Inches(5.3), y + Inches(0.05), Inches(3.2), Inches(0.5),
                     symptom, font_size=13)
        add_text_box(slide, Inches(8.5), y + Inches(0.05), Inches(4), Inches(0.5),
                     fix, font_size=13, color=ACCENT_GREEN)

    # ── SLIDE 12: Dominant Dimension ──────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "Trade-Off Analysis: Dominant Dimension")

    # Dominant
    add_box_shape(slide, Inches(0.8), Inches(1.8), Inches(3.5), Inches(2),
                  "DOMINANT\n\nStreaming\n\nOdds movement tracking\nis the core value prop.",
                  ACCENT_BLUE, font_size=13, text_color=WHITE)

    # Sacrificed
    add_box_shape(slide, Inches(5), Inches(1.8), Inches(3.5), Inches(2),
                  "SACRIFICED\n\nGeo-Distribution\n\nSingle region (us-east-1).\nNo multi-region needed\nfor analytics.",
                  RGBColor(200, 80, 80), font_size=13, text_color=WHITE)

    add_box_shape(slide, Inches(9.2), Inches(1.8), Inches(3.5), Inches(2),
                  "SACRIFICED\n\nFan-Out\n\nSingle producer, few\nconsumers. Pull-based\nAthena queries.",
                  RGBColor(200, 80, 80), font_size=13, text_color=WHITE)

    # Preserved
    add_text_box(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
                 "Preserved (but not dominant):", font_size=18, bold=True, color=NAVY)

    add_box_shape(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.2),
                  "CAP: Eventual consistency accepted\n\nGold tables may be briefly stale during\na pipeline run. Analysts query completed snapshots.",
                  WHITE, font_size=12)

    add_box_shape(slide, Inches(7), Inches(5.0), Inches(5.5), Inches(1.2),
                  "Caching: Not required today\n\nAthena scans are fast and cheap at current\nvolume. Redis layer is a natural future addition.",
                  WHITE, font_size=12)

    # ── SLIDE 13: Demo ────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
                 "Live Demo", font_size=48, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(2),
                 "1. Run the batch pipeline (ingest + transform + upload)\n"
                 "2. dbt run  -  materialize star schema in Athena\n"
                 "3. dbt test  -  validate 45 data quality checks\n"
                 "4. Query fact_game_odds in the Athena console",
                 font_size=20, color=RGBColor(180, 200, 230),
                 alignment=PP_ALIGN.CENTER)

    # ── SLIDE 14: Q&A ─────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
                 "Questions?", font_size=54, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
                 "github.com/daffleck7/oddswatch", font_size=20,
                 color=RGBColor(160, 180, 210), alignment=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────
    output_path = "docs/Team2_FinalPresentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")


if __name__ == "__main__":
    build_presentation()
